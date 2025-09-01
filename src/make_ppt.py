from pptx import Presentation
import toml

config = toml.load("config.toml")
OUT_PATH = config["presentation"]["out_path"]


def generate_presentation(themes):
    if not themes:
        print("No themes to generate PPT")
        return

    prs = Presentation()
    for theme in themes:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = theme
        slide.placeholders[1].text = "Articles for this theme go here.."

    prs.save(OUT_PATH)
    print(f"PPT saved to {OUT_PATH}")
